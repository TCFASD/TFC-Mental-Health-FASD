# -*- coding: utf-8 -*-
"""Build Module 1 draft deck on the Capacity-vs-Compliance theme. Open Sans, matched sizing, TFC logo, presenter."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = r"\\tfcfs2\Home\tamra.cajo\My Documents\Power-point\Capacity versus Compliance\Capacity versus Compliance in Court involved Systems.pptx"
OUT = r"\\tfcfs2\Home\tamra.cajo\My Documents\TFC-Mental-Health-FASD\Training\Module_1_Welcome_Why_Brain-Based_DRAFT.pptx"
LOGO = r"\\tfcfs2\Home\tamra.cajo\My Documents\My Pictures\logo no background.png"

NAVY  = RGBColor(0x34,0x40,0x68)
TEAL  = RGBColor(0x14,0x82,0xAB)
CYAN  = RGBColor(0x28,0xC4,0xCC)
GTEAL = RGBColor(0x42,0xBA,0x97)
GREEN = RGBColor(0x3E,0x88,0x53)
MUTED = RGBColor(0x62,0xA3,0x9F)
LIGHT = RGBColor(0xD9,0xE0,0xE6)
WHITE = RGBColor(0xFF,0xFF,0xFF)
TEXT  = RGBColor(0x33,0x37,0x3D)
F = "Open Sans"   # all wording in Open Sans (matches the Compliance deck)

prs = Presentation(TEMPLATE)
SW, SH = prs.slide_width/914400, prs.slide_height/914400

sldIdLst = prs.slides._sldIdLst
for sld in list(sldIdLst):
    rId = sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId); sldIdLst.remove(sld)
blank = min(prs.slide_layouts, key=lambda L: len(L.placeholders))

def slide(): return prs.slides.add_slide(blank)
def rect(s,x,y,w,h,color,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False
    return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=6,line_sp=1.08):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.space_before=Pt(0); p.line_spacing=line_sp
        for (t,sz,col,bold,ital) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.color.rgb=col
            r.font.bold=bold; r.font.italic=ital; r.font.name=F
    return tb
def bg(s,color=WHITE): rect(s,-0.1,-0.1,SW+0.2,SH+0.2,color)
def sidebar(s): rect(s,0,0,0.16,SH,NAVY); rect(s,0.16,0,0.06,SH,CYAN)
def title(s,t):
    txt(s,0.6,0.45,SW-2.6,1.0,[[(t,36,NAVY,False,False)]])
    rect(s,0.62,1.32,1.5,0.06,CYAN)
def logo_corner(s):
    s.shapes.add_picture(LOGO,Inches(SW-2.05),Inches(SH-0.92),width=Inches(1.7))
def R(t,sz,col=TEXT,bold=False,ital=False): return (t,sz,col,bold,ital)

# 1 — TITLE
s=slide(); bg(s,NAVY)
rect(s,0,4.55,SW,0.10,TEAL); rect(s,0,4.65,SW,0.10,GTEAL); rect(s,0,4.75,SW,0.10,CYAN)
txt(s,0.9,1.55,SW-1.8,2.2,[[R("Welcome to the FASD",44,WHITE)],[R("Mental Health Team",44,WHITE)]],sp_after=2)
txt(s,0.92,3.35,SW-1.8,0.7,[[R("Module 1  —  Why This Work Matters",24,CYAN)]])
txt(s,0.92,5.12,7.0,1.0,[[R("Brain-Based Mental Health Training",16,LIGHT)],
    [R("The Florida Center for Early Childhood",16,LIGHT)],
    [R("Presented by Tamra Cajo, LCSW",16,LIGHT)]],sp_after=1,line_sp=1.1)
# logo on white chip (navy text would vanish on navy bg)
rect(s,9.7,5.25,2.95,1.55,WHITE,MSO_SHAPE.ROUNDED_RECTANGLE)
s.shapes.add_picture(LOGO,Inches(9.95),Inches(5.5),width=Inches(2.45))

# 2 — WELCOME
s=slide(); bg(s); sidebar(s); title(s,"Welcome!")
txt(s,0.62,1.8,SW-1.4,4.4,
    [[R("You're joining the ",22),R("FASD Mental Health Division",22,TEAL,True),R(" at The Florida Center for Early Childhood.",22)],
     [R("",10)],
     [R("We support children and families affected by FASD and other ",22),R("brain-based (neurobehavioral) conditions.",22,TEAL,True)],
     [R("",10)],
     [R("Whether you're an intern, a new hire, or a licensed therapist — this training gives you what you need to start strong.",22)]],
    line_sp=1.18,sp_after=4)
logo_corner(s)

# 3 — BIG IDEA
s=slide(); bg(s,NAVY); rect(s,0,0,0.16,SH,CYAN)
txt(s,1.0,1.35,SW-2.0,0.9,[[R("The Big Idea",28,CYAN)]])
txt(s,1.0,2.35,SW-1.6,1.8,
    [[R("“What looks like ",40,WHITE),R("won't",40,CYAN,True,True),R(" is often ",40,WHITE),R("can't.",40,GTEAL,True,True),R("”",40,WHITE)]],line_sp=1.12)
txt(s,1.02,4.55,SW-2.0,1.6,
    [[R("Behavior is brain-based. Memory, attention, processing speed, focus, and language all shape how a child shows up — it's not about willpower.",22,LIGHT)]],line_sp=1.2)

# 4 — WHY (3 cards)
s=slide(); bg(s); sidebar(s); title(s,"Why a Brain-Based Lens?")
cards=[(TEAL,"No diagnosis needed","You don't need a formal FASD diagnosis to help. The lens works for any brain-based difference."),
       (GTEAL,"From shame to understanding","When families understand the brain, blame goes down — and symptoms often ease."),
       (GREEN,"Build on strengths","Every brain is different. We personalize care and start from what's already strong.")]
cw=3.85; gap=0.38; x0=0.62; y0=1.8; ch=3.95
for i,(col,h,b) in enumerate(cards):
    x=x0+i*(cw+gap)
    rect(s,x,y0,cw,ch,col,MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s,x+0.35,y0+0.45,0.7,0.09,WHITE)
    txt(s,x+0.32,y0+0.7,cw-0.64,1.3,[[R(h,20,WHITE,True)]],line_sp=1.0)
    txt(s,x+0.32,y0+1.95,cw-0.64,ch-2.1,[[R(b,16,WHITE)]],line_sp=1.18)
logo_corner(s)

# 5 — WHOLE SYSTEM
s=slide(); bg(s); sidebar(s); title(s,"We Support the Whole System")
txt(s,0.62,1.75,SW-1.4,1.3,[[R("The child, caregivers, and other supports may each need ",22),R("different accommodations",22,TEAL,True),R(".",22)]],line_sp=1.18)
txt(s,0.62,2.95,SW-1.4,0.5,[[R("We work as part of a team — we refer and coordinate for:",20)]])
items=["Further testing / evaluations","Occupational therapy","Speech & language","Educational navigation"]
y=3.7
for i,it in enumerate(items):
    rect(s,0.78,y+0.06,0.24,0.24,[TEAL,GTEAL,GREEN,MUTED][i],MSO_SHAPE.OVAL)
    txt(s,1.22,y,7.0,0.5,[[R(it,20,TEXT)]]); y+=0.62
txt(s,0.62,6.45,SW-2.4,0.6,[[R("…and we stay in contact with the child's other providers.",18,MUTED,False,True)]])
logo_corner(s)

# 6 — FAMILIES MAY NEED MORE
s=slide(); bg(s); sidebar(s); title(s,"FASD Families May Need More")
txt(s,0.62,1.5,SW-1.4,0.5,[[R("Supports inside our program and out in the community:",18,MUTED,False,True)]])
supports=[(TEAL,"Educational Navigation","Caregiver experts with lived experience help families advocate at school (e.g., Wright's Law)."),
          (GTEAL,"LRT — Local Review Team","Brings multiple agencies together when a child needs more than one program can provide."),
          (GREEN,"Parent Advocate","A dedicated advocate walking alongside the family."),
          (MUTED,"APD","Agency for Persons with Disabilities supports."),
          (CYAN,"CAT Team","Community Action Team support."),
          (TEAL,"TFC Supports","Other therapies and services here at The Florida Center.")]
cw=3.85; gap=0.38; ch=1.85; x0=0.62; y0=2.15
for i,(col,h,b) in enumerate(supports):
    r=i//3; c=i%3; x=x0+c*(cw+gap); y=y0+r*(ch+0.28)
    rect(s,x,y,0.10,ch,col); rect(s,x+0.10,y,cw-0.10,ch,LIGHT,MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s,x+0.34,y+0.2,cw-0.55,0.6,[[R(h,18,NAVY,True)]],line_sp=1.0)
    txt(s,x+0.34,y+0.78,cw-0.55,ch-0.9,[[R(b,14,TEXT)]],line_sp=1.12)
logo_corner(s)

# 7 — WHERE WE WORK
s=slide(); bg(s); sidebar(s); title(s,"Where We Work")
txt(s,0.62,2.0,8.6,3.0,
    [[R("The Florida Center is a ",22),R("Florida-licensed agency",22,TEAL,True),R(" — so we can serve Medicaid clients ",22),R("anywhere in the state.",22,TEAL,True)],
     [R("",12)],
     [R("There are no county or regional limits. These services are delivered by ",22),R("telehealth (virtually).",22,GTEAL,True)]],line_sp=1.2,sp_after=4)
rect(s,9.5,2.1,3.1,3.1,LIGHT,MSO_SHAPE.OVAL)
txt(s,9.5,3.15,3.1,1.1,[[R("Statewide",24,NAVY,True)],[R("Telehealth",20,TEAL)]],align=PP_ALIGN.CENTER,sp_after=2)
logo_corner(s)

# 8 — WHAT THIS MEANS FOR YOU
s=slide(); bg(s); sidebar(s); title(s,"What This Means for You")
pts=[("Bring the lens with you","You don't need to become an FASD specialist — add the brain-based lens to the work you already do."),
     ("Look underneath the behavior","Target the underlying skills (memory, focus, regulation), not just what's on the surface."),
     ("Stay reflective","Keep checking our own language, assumptions, and expectations.")]
y=1.85
for i,(h,b) in enumerate(pts):
    rect(s,0.7,y,0.55,0.55,[TEAL,GTEAL,GREEN][i],MSO_SHAPE.OVAL)
    txt(s,0.7,y,0.55,0.55,[[R(str(i+1),22,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,1.5,y-0.06,SW-2.4,1.4,[[R(h,21,NAVY,True)],[R(b,18,TEXT)]],line_sp=1.12,sp_after=2)
    y+=1.45
txt(s,1.5,y+0.05,SW-2.4,0.7,[[R("Let's get started.",24,TEAL,True,True)]])
logo_corner(s)

prs.save(OUT)
print("Saved:",OUT,"| slides:",len(prs.slides._sldIdLst))
