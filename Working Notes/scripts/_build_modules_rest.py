# -*- coding: utf-8 -*-
"""Build Modules 2,4,5,6,8 on the Capacity-vs-Compliance theme. Open Sans, TFC logo, matched sizing."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = r"\\tfcfs2\Home\tamra.cajo\My Documents\Power-point\Capacity versus Compliance\Capacity versus Compliance in Court involved Systems.pptx"
OUTDIR = r"\\tfcfs2\Home\tamra.cajo\My Documents\TFC-Mental-Health-FASD\Training"
LOGO = r"\\tfcfs2\Home\tamra.cajo\My Documents\My Pictures\logo no background.png"

NAVY=RGBColor(0x34,0x40,0x68); TEAL=RGBColor(0x14,0x82,0xAB); CYAN=RGBColor(0x28,0xC4,0xCC)
GTEAL=RGBColor(0x42,0xBA,0x97); GREEN=RGBColor(0x3E,0x88,0x53); MUTED=RGBColor(0x62,0xA3,0x9F)
LIGHT=RGBColor(0xD9,0xE0,0xE6); WHITE=RGBColor(0xFF,0xFF,0xFF); TEXT=RGBColor(0x33,0x37,0x3D)
ACC=[TEAL,GTEAL,GREEN,MUTED,CYAN]
F="Open Sans"
PRS=None; BLANK=None; SW=0; SH=0

def new_deck():
    global PRS,BLANK,SW,SH
    PRS=Presentation(TEMPLATE); SW=PRS.slide_width/914400; SH=PRS.slide_height/914400
    lst=PRS.slides._sldIdLst
    for sld in list(lst):
        rId=sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        PRS.part.drop_rel(rId); lst.remove(sld)
    BLANK=min(PRS.slide_layouts,key=lambda L:len(L.placeholders))

def slide(): return PRS.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False
    return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=6,line_sp=1.08):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.space_before=Pt(0); p.line_spacing=line_sp
        for (t,sz,col,bold,ital) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.color.rgb=col
            r.font.bold=bold; r.font.italic=ital; r.font.name=F
    return tb
def bg(s,color=WHITE): rect(s,-0.1,-0.1,SW+0.2,SH+0.2,color)
def sidebar(s): rect(s,0,0,0.16,SH,NAVY); rect(s,0.16,0,0.06,SH,CYAN)
def title(s,t,size=36):
    txt(s,0.6,0.42,SW-2.5,1.05,[[(t,size,NAVY,False,False)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,0.62,1.34,1.5,0.06,CYAN)
def logo_corner(s): s.shapes.add_picture(LOGO,Inches(SW-2.05),Inches(SH-0.92),width=Inches(1.7))
def R(t,sz,col=TEXT,bold=False,ital=False): return (t,sz,col,bold,ital)

def title_slide(main_lines, module_sub):
    s=slide(); bg(s,NAVY)
    rect(s,0,4.55,SW,0.10,TEAL); rect(s,0,4.65,SW,0.10,GTEAL); rect(s,0,4.75,SW,0.10,CYAN)
    runs=[[R(l,44,WHITE)] for l in main_lines]
    txt(s,0.9,1.5,SW-1.8,2.0,runs,sp_after=2)
    txt(s,0.92,3.5,SW-1.8,0.7,[[R(module_sub,24,CYAN)]])
    txt(s,0.92,5.12,7.0,1.0,[[R("Brain-Based Mental Health Training",16,LIGHT)],
        [R("The Florida Center for Early Childhood",16,LIGHT)],
        [R("Presented by Tamra Cajo, LCSW",16,LIGHT)]],sp_after=1,line_sp=1.1)
    rect(s,9.7,5.25,2.95,1.55,WHITE,MSO_SHAPE.ROUNDED_RECTANGLE)
    s.shapes.add_picture(LOGO,Inches(9.95),Inches(5.5),width=Inches(2.45))
    return s

def navy_quote(label, quote_runs, sub):
    s=slide(); bg(s,NAVY); rect(s,0,0,0.16,SH,CYAN)
    txt(s,1.0,1.35,SW-2.0,0.9,[[R(label,28,CYAN)]])
    txt(s,1.0,2.35,SW-1.6,1.8,[quote_runs],line_sp=1.12)
    txt(s,1.02,4.55,SW-2.0,1.6,[[R(sub,22,LIGHT)]],line_sp=1.2)
    return s

def cards_row(s, cards, y0=1.8, ch=3.95, body_sz=16, head_sz=20):
    n=len(cards); cw=(SW-1.24-0.38*(n-1))/n; x0=0.62
    for i,(col,h,b) in enumerate(cards):
        x=x0+i*(cw+0.38)
        rect(s,x,y0,cw,ch,col,MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s,x+0.32,y0+0.42,0.7,0.09,WHITE)
        txt(s,x+0.3,y0+0.66,cw-0.6,1.3,[[R(h,head_sz,WHITE,True)]],line_sp=1.0)
        txt(s,x+0.3,y0+1.9,cw-0.6,ch-2.05,[[R(b,body_sz,WHITE)]],line_sp=1.16)

def bullets(s, items, x=0.7, y=1.85, w=None, gap=0.7, sz=20, dot=True):
    if w is None: w=SW-x-0.6
    for i,it in enumerate(items):
        yy=y+i*gap
        if dot: rect(s,x,yy+0.07,0.24,0.24,ACC[i%len(ACC)],MSO_SHAPE.OVAL)
        off=0.42 if dot else 0
        if isinstance(it,tuple):
            txt(s,x+off,yy-0.04,w-off,gap,[[R(it[0]+"  ",sz,NAVY,True),R(it[1],sz-2,TEXT)]],line_sp=1.1)
        else:
            txt(s,x+off,yy,w-off,gap,[[R(it,sz,TEXT)]],line_sp=1.1)

def softcards(s, items, y0=1.7, ch=1.85, cols=3, head_sz=18, body_sz=14):
    cw=(SW-1.24-0.38*(cols-1))/cols; x0=0.62
    for i,(col,h,b) in enumerate(items):
        r=i//cols; c=i%cols; x=x0+c*(cw+0.38); y=y0+r*(ch+0.28)
        rect(s,x,y,0.10,ch,col); rect(s,x+0.10,y,cw-0.10,ch,LIGHT,MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s,x+0.32,y+0.18,cw-0.52,0.7,[[R(h,head_sz,NAVY,True)]],line_sp=1.0)
        txt(s,x+0.32,y+0.74,cw-0.52,ch-0.85,[[R(b,body_sz,TEXT)]],line_sp=1.12)

def lead(s,y,runs,sz=22): txt(s,0.62,y,SW-1.4,1.2,[runs],line_sp=1.18)

# ============================================================ MODULE 2
def build_m2():
    new_deck()
    title_slide(["Understanding","the Brain"],"Module 2  —  The Brain-Based Lens")
    # 2 levels
    s=slide(); bg(s); sidebar(s); title(s,"Two Ways to Practice")
    cards_row(s,[(TEAL,"Brain-Based Informed","The foundation. You understand how the brain shapes behavior, learning, and emotion — and you're beginning to apply it."),
                 (GREEN,"Brain-Based Responsive","The next level. You read the moment and adapt in real time — pacing, environment, and approach shift with the person.")],ch=3.4,body_sz=18,head_sz=22)
    logo_corner(s)
    # core competencies
    s=slide(); bg(s); sidebar(s); title(s,"What the Lens Covers")
    softcards(s,[(TEAL,"Brain & Behavior","How brain development and neuroplasticity shape what we see."),
                 (GTEAL,"Trauma + the Brain","How trauma changes regulation, memory, and stress response."),
                 (GREEN,"Prenatal Alcohol Exposure","How PAE affects the brain — across every age."),
                 (MUTED,"Cause & Effect","Why abstract thinking and consequences work differently."),
                 (CYAN,"Client Education","Explaining the brain in ways that reduce shame."),
                 (TEAL,"Case Conceptualization","Weaving brain-based factors into the whole plan.")],y0=1.75,ch=1.9)
    logo_corner(s)
    # trauma
    s=slide(); bg(s); sidebar(s); title(s,"Trauma + the Brain")
    lead(s,1.8,[R("Trauma reshapes the brain areas that handle ",22),R("emotion, memory, stress, and self-control.",22,TEAL,True)])
    bullets(s,[("Window of tolerance","each person has a zone where they can think and connect."),
               ("Safety first","the brain scans for safety or threat before anything else."),
               ("Bottom-up","calm the body and build regulation before the thinking work.")],y=3.0,sz=19)
    logo_corner(s)
    # PAE
    s=slide(); bg(s); sidebar(s); title(s,"Prenatal Alcohol Exposure (PAE)")
    lead(s,1.75,[R("PAE creates ",22),R("organic brain differences",22,TEAL,True),R(" — in self-regulation, memory, focus, and social skills.",22)])
    bullets(s,[("Screen at every age","many people are never diagnosed."),
               ("A negative screen isn't a clear","it doesn't rule FASD out."),
               ("Build accommodations","and bring the support system into treatment.")],y=3.0,sz=19)
    logo_corner(s)
    # cause-effect
    s=slide(); bg(s); sidebar(s); title(s,"Cause & Effect Works Differently")
    lead(s,1.8,[R("Brain differences can make ",22),R("cause-and-effect reasoning",22,TEAL,True),R(" genuinely hard.",22)])
    bullets(s,[("Delayed consequences often don't land","the connection isn't made."),
               ("Keep it immediate and concrete","and teach the link out loud."),
               ("Meet the brain where it is","adjust to their capacity in the moment.")],y=3.0,sz=19)
    logo_corner(s)
    # language
    s=slide(); bg(s); sidebar(s); title(s,"Words Matter")
    lead(s,1.9,[R("The ",22),R("Language & Stigma Guide",22,TEAL,True),R(" is required reading before you see clients.",22)])
    txt(s,0.62,3.0,SW-1.4,2.0,[[R("The words we choose can add shame — or take it away. Brain-based language describes what's happening in the brain instead of blaming the person for it.",21,TEXT)]],line_sp=1.25)
    logo_corner(s)
    # required watch
    s=slide(); bg(s); sidebar(s); title(s,"Your Required Watch")
    softcards(s,[(TEAL,"Screening for Prenatal Alcohol Exposure","Zoom recording  ·  passcode SCREEN#02242025"),
                 (GREEN,"MH Through a Brain-Based Lens","Zoom recording  ·  passcode TFCMENTALHEALTH#2025"),
                 (MUTED,"Baseline Self-Assessment","Complete the 50-item tool BEFORE you watch — you'll re-take it later.")],y0=1.9,ch=2.3,cols=3,body_sz=15)
    txt(s,0.62,4.6,SW-1.4,0.6,[[R("Watch these first — the live sessions build on them.",18,MUTED,False,True)]])
    logo_corner(s)
    PRS.save(OUTDIR+r"\Module_2_Understanding_the_Brain_DRAFT.pptx")

# ============================================================ MODULE 4
def build_m4():
    new_deck()
    title_slide(["Assessment"],"Module 4  —  Through a Brain-Based Lens")
    s=slide(); bg(s); sidebar(s); title(s,"Why Assess Brain Function?")
    lead(s,1.8,[R("Traditional mental-health assessments often ",22),R("miss the brain-based challenges",22,TEAL,True),R(" driving behavior and emotion.",22)])
    txt(s,0.62,3.0,SW-1.4,0.6,[[R("Get the fuller picture — pull in what others already know:",20)]])
    bullets(s,["Neuropsychological evaluations","Occupational & speech-language evals","School reports and IEPs","Medical and medication history"],y=3.8,sz=19)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Before the First Session")
    lead(s,1.8,[R("Gather what exists and listen to the story — including ",22),R("prenatal and birth history.",22,TEAL,True)])
    bullets(s,[("Request records early","prior evals, IEPs, medical history."),
               ("Screen for PAE","ask with stigma-reducing language."),
               ("Note strengths too","not just the challenges.")],y=3.0,sz=19)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Which Assessment?")
    softcards(s,[(TEAL,"IDA","Children 0–5 (In-Depth Assessment)."),
                 (GTEAL,"BPS — Child","Children 6–18 (Biopsychosocial)."),
                 (GREEN,"BPS — Adult","Adult clients.")],y0=1.85,ch=1.7,cols=3,body_sz=15)
    txt(s,0.62,3.95,SW-1.4,1.4,[[R("Shortcut: an ",19),R("FASD evaluation within the last 6 months",19,TEAL,True),R(" can stand in for the BPS — go straight to the Treatment Plan. Older than 6 months? A BPS or IDA is needed first.",19)]],line_sp=1.2)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Required at Assessment")
    bullets(s,[("Child present","the child must be there for the assessment."),
               ("Brief Status Exam (BSE)","by a licensed practitioner before the Treatment Plan."),
               ("CFARS","at start and quarterly (interns can't bill it)."),
               ("Outcome measures","baseline now, every 3 months, and at closing.")],y=1.9,sz=19,gap=0.85)
    txt(s,1.12,5.5,SW-2.0,0.8,[[R("Outcomes: Parenting Sense of Competency · Perceived Stress Scale · Neurobehavioral Screener (caregiver + child).",16,MUTED,False,True)]],line_sp=1.15)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"When There's No Neuropsych")
    lead(s,1.8,[R("No formal evaluation? You can still ",22),R("read the brain informally.",22,TEAL,True)])
    cards_row(s,[(TEAL,"Ability Wheel","Vanessa Spiller's simple tool for cognitive and functional skills."),
                 (GREEN,"FASCETS Exploration Tool","A structured neurobehavioral tool — requires FASCETS training.")],y0=3.0,ch=2.5,body_sz=17,head_sz=21)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Map the Brain by Domain")
    softcards(s,[(TEAL,"Memory","Children's Memory Scale · NEPSY-II"),
                 (GTEAL,"Development","Bayley · Vineland"),
                 (GREEN,"Sensory","Sensory Profile"),
                 (MUTED,"Processing / IQ","WISC · WAIS"),
                 (CYAN,"Speech & Language","CASL-2"),
                 (TEAL,"Auditory Processing","SCAN")],y0=1.8,ch=1.55,body_sz=14)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Connect the Dots")
    txt(s,0.9,1.85,SW-1.8,1.0,[[R("Report  →  Behavior  →  Accommodation",26,NAVY,True)]],align=PP_ALIGN.CENTER)
    rect(s,3.5,2.7,SW-7.0,0.05,LIGHT)
    txt(s,0.62,3.3,SW-1.4,0.6,[[R("And don't stop at the child — also assess:",20)]])
    bullets(s,["The caregiver's learning style","Stress in the everyday environment","Resources for the client and their supports"],y=4.1,sz=19)
    logo_corner(s)
    PRS.save(OUTDIR+r"\Module_4_Assessment_Brain-Based_DRAFT.pptx")

# ============================================================ MODULE 5
def build_m5():
    new_deck()
    title_slide(["Treatment","Planning"],"Module 5  —  Through a Brain-Based Lens")
    s=slide(); bg(s); sidebar(s); title(s,"Planning Is Collaborative")
    lead(s,1.85,[R("The Treatment Plan is ",22),R("required before services",22,TEAL,True),R(" — and it's built together.",22)])
    txt(s,0.62,3.0,SW-1.4,0.6,[[R("Around the table:",20)]])
    bullets(s,["The family (and the client, where appropriate)","The clinician and supervisor","Other treating practitioners","Anyone the family wants included"],y=3.8,sz=19)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"The Planning Conversation")
    txt(s,0.62,1.6,SW-1.4,0.6,[[R("Start with the outcome measures, then ask:",20,TEXT)]])
    bullets(s,["What do you hope to get out of treatment?","What's the goal for the parent? For the child?","How does meeting on a regular schedule feel?","Is this realistic for the child's brain — and where the caregiver is right now?"],y=2.5,sz=19,gap=0.82)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Brain-Based Planning")
    cards_row(s,[(TEAL,"Include the supports","Caregivers and key adults shape the environment that drives symptoms."),
                 (GTEAL,"Aim for interdependence","Healthy reliance on others — not independence at all costs."),
                 (GREEN,"Build trust first","Felt safety comes before the harder work.")],ch=3.6,body_sz=16,head_sz=20)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Strengths-Based Goals")
    bullets(s,[("Developmental age, not birthday","match goals to where the brain is."),
               ("Small, concrete steps","with repetition and flexible pacing."),
               ("Consistency everywhere","home, school, and therapy together."),
               ("Shape the environment","lighting, clutter, and transitions matter.")],y=1.9,sz=19,gap=0.9)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Required Content & Signatures")
    softcards(s,[(TEAL,"Measurable objectives","Specific and trackable — not other services."),
                 (GTEAL,"Crisis / Safety Plan","Always included."),
                 (GREEN,"Discharge Plan","Restates the measurable objectives."),
                 (MUTED,"Service table","List every service (e.g., add Group)."),
                 (CYAN,"Give parents a copy","Every time."),
                 (TEAL,"Signatures","Supervisor + 'Type 05' for Medicaid.")],y0=1.8,ch=1.7,body_sz=14)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Deadlines That Matter")
    bullets(s,[("Within 10 calendar days","complete the plan after the assessment."),
               ("Every 90 days","review to keep the family eligible and progressing."),
               ("Changing mid-stream?","add an Addendum — treatment follows the latest plan.")],y=2.0,sz=20,gap=1.0)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Getting It Signed (FOXIT)")
    bullets(s,["Save the completed plan with therapist/supervisor signature","Upload to FOXIT and add the caregiver's email","Drop in the signature + date fields, then Send","Save the signed plan back into Lauris, labeled correctly"],y=2.0,sz=20,gap=1.0,dot=True)
    logo_corner(s)
    PRS.save(OUTDIR+r"\Module_5_Treatment_Planning_Brain-Based_DRAFT.pptx")

# ============================================================ MODULE 6
def build_m6():
    new_deck()
    title_slide(["Treatment &","Adapting"],"Module 6  —  Your Approach")
    s=slide(); bg(s); sidebar(s); title(s,"Evidence-Based, Brain-Adapted")
    softcards(s,[(TEAL,"Parent Management (PMT)","PCIT and Triple P — strengthen the parent-child relationship."),
                 (GTEAL,"CBT","Reshape unhelpful thoughts and build coping skills."),
                 (GREEN,"DBT","Mindfulness, emotion regulation, distress tolerance."),
                 (MUTED,"TBRI","Trauma-informed: connection, empowerment, correction.")],y0=1.9,ch=2.3,cols=2,head_sz=19,body_sz=15)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Adjust for the Brain")
    softcards(s,[(TEAL,"Low IQ / younger dev. age","Slow down, simplify, and repeat across sessions."),
                 (GTEAL,"Memory or carry-over trouble","One skill at a time · caregivers practice daily · video skills to review."),
                 (GREEN,"Concrete, literal thinking","One rule at a time · role-play · frequent review.")],y0=1.9,ch=2.4,cols=3,body_sz=14.5)
    txt(s,0.62,4.55,SW-1.4,0.6,[[R("Same evidence-based treatment — paced and shaped to how the brain works.",18,MUTED,False,True)]])
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Documentation: DAP Notes")
    lead(s,1.7,[R("Every session, in ",22),R("DAP format",22,TEAL,True),R(" — Data, Assessment, Plan.",22)])
    bullets(s,[("Data","setting, techniques used, and the child's response."),
               ("Assessment","the mental status exam (mood, affect, cognition)."),
               ("Plan","next steps, with the next appointment date and time."),
               ("Within 48 hours","and document every no-show or cancellation.")],y=2.9,sz=18,gap=0.82)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Closing a Case")
    bullets(s,[("Finish all notes","make sure progress notes are complete."),
               ("Discharge Plan for everyone","with recommended follow-up services."),
               ("Close in Lauris","complete the Close Form; supervisor signs off."),
               ("Aim for the 1st","try to close clients before the start of the month.")],y=1.95,sz=19,gap=0.92)
    logo_corner(s)
    PRS.save(OUTDIR+r"\Module_6_Treatment_and_Adapting_DRAFT.pptx")

# ============================================================ MODULE 8
def build_m8():
    new_deck()
    title_slide(["Billing &","Productivity"],"Module 8  —  The Essentials")
    s=slide(); bg(s); sidebar(s); title(s,"The Billing Absolutes")
    items=[("1","No client is seen until consents are signed."),
           ("2","The assessment is always billed first."),
           ("3","Treatment Plan within 10 calendar days — caregiver signs."),
           ("4","Track expiration dates — renew before you bill again.")]
    y=1.95
    for i,(n,b) in enumerate(items):
        rect(s,0.7,y,0.6,0.6,ACC[i],MSO_SHAPE.OVAL)
        txt(s,0.7,y,0.6,0.6,[[R(n,24,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,1.55,y+0.02,SW-2.4,0.7,[[R(b,21,TEXT)]],anchor=MSO_ANCHOR.MIDDLE,line_sp=1.05)
        y+=1.0
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Service Codes at a Glance")
    softcards(s,[(TEAL,"Assessments","PSY / BPS · IDA · CFARS — billed per event."),
                 (GTEAL,"Treatment Plan","H0032 — plus reviews, per event."),
                 (GREEN,"Therapy","Individual, family, or group — billed in time units."),
                 (MUTED,"Consult / CCM","Phone & collateral — not billable to insurance.")],y0=1.85,ch=2.0,cols=2,head_sz=19,body_sz=15)
    txt(s,0.62,4.3,SW-1.4,0.8,[[R("Always pull the client up in Lauris first — confirm insurance and the available code before you bill. Phone-only is never reimbursable.",17,MUTED,False,True)]],line_sp=1.2)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Who Can Bill What")
    bullets(s,[("Child + caregiver present","no child in the room, no billable service."),
               ("Diagnosis is signed by a licensed clinician","others can collect the data."),
               ("BPS vs. IDA","bachelor-level/interns can bill a BPS, not an IDA."),
               ("Hand-off matters","a master's clinician completes the TP within 10 days.")],y=1.95,sz=18,gap=0.92)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"The Timeline")
    bullets(s,[("Assessment","completed within 5 business days of the first visit."),
               ("Treatment Plan","within 10 calendar days of the assessment."),
               ("Reviews","every 90 days keeps services going."),
               ("Diagnosis changes","only at a Review or Assessment.")],y=1.95,sz=19,gap=0.92)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Child vs. Adult  &  Interns")
    cards_row(s,[(TEAL,"Billing the child or adult","You can bill under either — if you bill the adult, keep the focus on parenting."),
                 (GREEN,"Intern codes (the 'IN' suffix)","Non-billable, but the documentation standard is identical to a licensed clinician's.")],y0=2.0,ch=2.6,body_sz=17,head_sz=20)
    logo_corner(s)
    s=slide(); bg(s); sidebar(s); title(s,"Productivity: the 75 / 25")
    txt(s,0.9,1.8,SW-1.8,1.0,[[R("75%",54,TEAL,True),R("  direct / billable      ",24,NAVY),R("25%",54,GTEAL,True),R("  admin",24,NAVY)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    bullets(s,[("About 2 hours a day","is administrative time."),
               ("Keep consultation small","around one hour a week."),
               ("Track it yourself","and raise caseload concerns early — supervisors won't guess.")],y=3.4,sz=19,gap=0.9)
    logo_corner(s)
    PRS.save(OUTDIR+r"\Module_8_Billing_and_Productivity_DRAFT.pptx")

for b in (build_m2,build_m4,build_m5,build_m6,build_m8):
    b(); print("built", b.__name__)
print("DONE")
